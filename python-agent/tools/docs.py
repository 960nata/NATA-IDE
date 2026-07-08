"""docs.py — tool buat bikin dokumen: Word, PDF, PowerPoint."""

from pathlib import Path

WORKSPACE = Path(__file__).resolve().parent.parent / "workspace"
WORKSPACE.mkdir(exist_ok=True)


def _safe(name, ext):
    """Pastiin nama file aman & ada ekstensinya."""
    name = Path(name).name  # buang path, ambil nama doang
    if not name.lower().endswith(ext):
        name += ext
    return WORKSPACE / name


def create_word(filename="dokumen.docx", title="", content=None):
    from docx import Document

    doc = Document()
    if title:
        doc.add_heading(title, level=0)
    for para in (content or []):
        doc.add_paragraph(str(para))
    path = _safe(filename, ".docx")
    doc.save(path)
    return f"✅ File Word dibuat: {path}"


def create_pdf(filename="dokumen.pdf", title="", content=None):
    from fpdf import FPDF

    pdf = FPDF()
    pdf.add_page()
    if title:
        pdf.set_font("Helvetica", "B", 16)
        pdf.multi_cell(0, 10, title)
        pdf.ln(4)
    pdf.set_font("Helvetica", size=12)
    for para in (content or []):
        # font inti fpdf cuma latin-1, buang karakter aneh biar ga error
        text = str(para).encode("latin-1", "replace").decode("latin-1")
        pdf.multi_cell(0, 8, text)
        pdf.ln(2)
    path = _safe(filename, ".pdf")
    pdf.output(str(path))
    return f"✅ File PDF dibuat: {path}"


def create_pptx(filename="presentasi.pptx", title="", slides=None):
    from pptx import Presentation

    prs = Presentation()

    # slide judul
    if title:
        layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(layout)
        slide.shapes.title.text = title

    # slide isi: tiap item = {"title": "...", "bullets": ["...", "..."]}
    for item in (slides or []):
        layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(layout)
        slide.shapes.title.text = item.get("title", "")
        body = slide.placeholders[1].text_frame
        bullets = item.get("bullets", [])
        if bullets:
            body.text = str(bullets[0])
            for b in bullets[1:]:
                body.add_paragraph().text = str(b)
    path = _safe(filename, ".pptx")
    prs.save(path)
    return f"✅ File PowerPoint dibuat: {path}"


# Daftar tool yang diekspos ke agent
TOOLS = [
    {
        "name": "create_word",
        "description": "Bikin dokumen Word (.docx).",
        "params": '{"filename": "nama.docx", "title": "judul", "content": ["paragraf 1", "paragraf 2"]}',
        "run": create_word,
    },
    {
        "name": "create_pdf",
        "description": "Bikin dokumen PDF (.pdf).",
        "params": '{"filename": "nama.pdf", "title": "judul", "content": ["paragraf 1", "paragraf 2"]}',
        "run": create_pdf,
    },
    {
        "name": "create_pptx",
        "description": "Bikin presentasi PowerPoint (.pptx).",
        "params": '{"filename": "nama.pptx", "title": "judul utama", "slides": [{"title": "judul slide", "bullets": ["poin 1", "poin 2"]}]}',
        "run": create_pptx,
    },
]
